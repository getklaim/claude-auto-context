#!/usr/bin/env python3
"""Auto Context PPT Generator — Mission Control Design System"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import os

# ─── Design Tokens ───────────────────────────────────────────
BG_BASE = RGBColor(0x0A, 0x0F, 0x1C)
BG_SURFACE1 = RGBColor(0x12, 0x1A, 0x2B)
BG_SURFACE2 = RGBColor(0x0F, 0x17, 0x27)
LINE_DEFAULT = RGBColor(0x26, 0x32, 0x4A)
TEXT_PRIMARY = RGBColor(0xE8, 0xEE, 0xF9)
TEXT_SECONDARY = RGBColor(0xA9, 0xB6, 0xCF)
TABLE_HEADER_BG = RGBColor(0x17, 0x22, 0x38)

STATE_OBSERVE = RGBColor(0x38, 0xBD, 0xF8)
STATE_CANDIDATE = RGBColor(0x22, 0xD3, 0xEE)
STATE_CONVENTION = RGBColor(0x22, 0xC5, 0x5E)
STATE_DECAY = RGBColor(0xF4, 0x3F, 0x5E)
STATE_APPROVAL = RGBColor(0xF5, 0x9E, 0x0B)

# ─── Typography Roles ─────────────────────────────────────────
# Title: Pretendard SemiBold (46~56)
# Body:  Pretendard Regular  (24~30)
# Data:  JetBrains Mono      (20~24)
FONT_TITLE = "Pretendard SemiBold"
FONT_BODY  = "Pretendard"
FONT_DATA  = "JetBrains Mono"

# Role-based font resolver
def _font(role):
    """role: 'title' | 'body' | 'data'"""
    return {
        "title": FONT_TITLE,
        "body":  FONT_BODY,
        "data":  FONT_DATA,
    }.get(role, FONT_BODY)

# px to EMU (1920x1080 canvas → 13333333 x 7500000 EMU)
SLIDE_W = 13333333
SLIDE_H = 7500000
PX_RATIO = SLIDE_W / 1920

def px(val):
    return int(val * PX_RATIO)

def add_bg(slide, color=BG_BASE):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color

def add_shape(slide, left, top, width, height, fill=None, border_color=None, border_width=Pt(1), radius=None):
    shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE if radius else MSO_SHAPE.RECTANGLE,
        px(left), px(top), px(width), px(height)
    )
    if fill:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill
    else:
        shape.fill.background()
    if border_color:
        shape.line.color.rgb = border_color
        shape.line.width = border_width
    else:
        shape.line.fill.background()
    return shape

def add_text(slide, left, top, width, height, text, font_size=24, color=TEXT_PRIMARY, bold=False, alignment=PP_ALIGN.LEFT, role="body"):
    """role: 'title' | 'body' | 'data' — determines font family"""
    txBox = slide.shapes.add_textbox(px(left), px(top), px(width), px(height))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = _font(role)
    p.alignment = alignment
    return txBox

def add_card(slide, left, top, width, height, title, body, accent_color, icon_text=None):
    shape = add_shape(slide, left, top, width, height, fill=BG_SURFACE1, border_color=accent_color)
    # title — Pretendard SemiBold
    add_text(slide, left+24, top+20, width-48, 40, title, font_size=20, color=accent_color, bold=True, role="title")
    # body — Pretendard Regular
    add_text(slide, left+24, top+72, width-48, height-92, body, font_size=16, color=TEXT_SECONDARY, role="body")
    return shape

def add_table_slide(slide, left, top, width, rows_data, col_widths_pct):
    """rows_data: list of lists. First row is header."""
    n_rows = len(rows_data)
    n_cols = len(rows_data[0])
    table_shape = slide.shapes.add_table(n_rows, n_cols, px(left), px(top), px(width), px(n_rows * 56))
    table = table_shape.table

    total_w = px(width)
    for ci, pct in enumerate(col_widths_pct):
        table.columns[ci].width = int(total_w * pct)

    for ri, row in enumerate(rows_data):
        for ci, cell_text in enumerate(row):
            cell = table.cell(ri, ci)
            cell.text = str(cell_text)
            for paragraph in cell.text_frame.paragraphs:
                paragraph.font.size = Pt(14)
                if ri == 0:
                    paragraph.font.name = FONT_TITLE
                    paragraph.font.bold = True
                    paragraph.font.color.rgb = TEXT_PRIMARY
                else:
                    paragraph.font.name = FONT_BODY
                    paragraph.font.color.rgb = TEXT_SECONDARY
            # cell fill
            cell_fill = cell.fill
            cell_fill.solid()
            if ri == 0:
                cell_fill.fore_color.rgb = TABLE_HEADER_BG
            elif ri % 2 == 0:
                cell_fill.fore_color.rgb = BG_SURFACE1
            else:
                cell_fill.fore_color.rgb = BG_BASE
    return table_shape


# ─── Create Presentation ─────────────────────────────────────
prs = Presentation()
prs.slide_width = SLIDE_W
prs.slide_height = SLIDE_H
blank_layout = prs.slide_layouts[6]  # blank


# ═══════════════════════════════════════════════════════════════
# SLIDE 1 — Title Hero
# ═══════════════════════════════════════════════════════════════
s1 = prs.slides.add_slide(blank_layout)
add_bg(s1)
# Grid accent line
add_shape(s1, 0, 0, 1920, 1080, fill=BG_BASE)
# Title
add_text(s1, 160, 300, 1600, 140,
    "Auto Context",
    font_size=56, color=TEXT_PRIMARY, bold=True, alignment=PP_ALIGN.CENTER, role="title")
# Subtitle
add_text(s1, 160, 460, 1600, 70,
    "Claude Code\uac00 \ucf54\ub529\ud560 \ub54c, \ud504\ub85c\uc81d\ud2b8\uac00 \uc2a4\uc2a4\ub85c \uc88b\uc740 \ucee8\ud14d\uc2a4\ud2b8\ub97c \ub9cc\ub4e4\uac8c \ud558\ub294 \ud50c\ub7ec\uadf8\uc778",
    font_size=24, color=TEXT_SECONDARY, alignment=PP_ALIGN.CENTER)
# Footer
add_text(s1, 160, 980, 400, 40,
    "Claude Code Plugin \xb7 2026",
    font_size=16, color=LINE_DEFAULT, alignment=PP_ALIGN.LEFT, role="data")


# ═══════════════════════════════════════════════════════════════
# SLIDE 2 — Pain: Context Window의 현실
# ═══════════════════════════════════════════════════════════════
s2 = prs.slides.add_slide(blank_layout)
add_bg(s2)
add_shape(s2, 0, 0, 1920, 1080, fill=BG_BASE)

add_text(s2, 96, 72, 1728, 100,
    "Context Window \u2014 AI\uac00 \ud55c \ubc88\uc5d0 \ubcfc \uc218 \uc788\ub294 \uc804\ubd80",
    font_size=34, color=TEXT_PRIMARY, bold=True, role="title")

add_text(s2, 96, 190, 1728, 90,
    "\u201cn\uac1c \ud1a0\ud070\uc740 n\xb2\uac1c\uc758 \uad00\uacc4\ub97c \uacc4\uc0b0\ud574\uc57c \ud55c\ub2e4. \ud1a0\ud070\uc774 2\ubc30\uba74 \uc5f0\uc0b0\uc740 4\ubc30. \uc8fc\uc758\ub825 \uc790\uccb4\uac00 \ud76c\uc11d\ub41c\ub2e4.\u201d",
    font_size=20, color=TEXT_SECONDARY, role="data")

# Card 1
add_card(s2, 96, 330, 544, 280,
    "Tool Output\uc774 83.9%",
    "\uc5d0\uc774\uc804\ud2b8 \uc791\uc5c5 \uc911 \ub3c4\uad6c \ucd9c\ub825\uc774\n\ucee8\ud14d\uc2a4\ud2b8\uc758 83.9%\ub97c \ucc28\uc9c0.\n\ub300\ubd80\ubd84\uc740 \uc774\ubbf8 \uc4f8\ubaa8\ub97c \ub2e4\ud55c \uc7a1\uc74c.",
    STATE_DECAY)

# Card 2
add_card(s2, 688, 330, 544, 280,
    "\ubb34\uad00 \ubb38\uc11c 1\uac1c = \uc131\ub2a5 \uae09\ub77d",
    "\uad00\ub828 \uc5c6\ub294 \ubb38\uc11c \ud558\ub098\ub9cc \ub123\uc5b4\ub3c4\nstep function\uc73c\ub85c \uc131\ub2a5 \ud558\ub77d.\n\ubaa8\ub378\uc740 \uc5b4\ub5a4 \uac83\ub3c4 \"\uac74\ub108\ub6f8\" \uc218 \uc5c6\ub2e4.",
    STATE_DECAY)

# Card 3
add_card(s2, 1280, 330, 544, 280,
    "32K+ \uc8fc\uc7a5 \uc911 50%\ub9cc \uc2e4\uc81c \uc720\uc9c0",
    "RULER \ubca4\uce58\ub9c8\ud06c \u2014 \uc7a5\ubb38 context\ub97c\n\uc8fc\uc7a5\ud558\ub294 \ubaa8\ub378 \uc911 \uc808\ubc18\ub9cc\n\uc2e4\uc81c \ud574\ub2f9 \uae38\uc774\uc5d0\uc11c \uc131\ub2a5 \uc720\uc9c0.",
    STATE_DECAY)

# Bottom message
add_text(s2, 96, 670, 1728, 120,
    "Context Window\uac00 \ud06c\ub2e4\uace0 \uc88b\uc740 \uac8c \uc544\ub2c8\ub2e4.\n\uc62c\ubc14\ub978 \uc815\ubcf4\ub9cc \ub4e4\uc5b4 \uc788\uc744 \ub54c \uc88b\uc740 \uac83\uc774\ub2e4.",
    font_size=24, color=STATE_OBSERVE, bold=True, alignment=PP_ALIGN.CENTER, role="title")


# ═══════════════════════════════════════════════════════════════
# SLIDE 3 — Benchmark: 모델별 열화 지점
# ═══════════════════════════════════════════════════════════════
s3 = prs.slides.add_slide(blank_layout)
add_bg(s3)
add_shape(s3, 0, 0, 1920, 1080, fill=BG_BASE)

add_text(s3, 96, 72, 1728, 90,
    "\ub354 \ub123\uc73c\uba74 \ub354 \ub098\ube60\uc9c4\ub2e4 \u2014 \ubaa8\ub378\ubcc4 \uc5f4\ud654 \uc9c0\uc810",
    font_size=34, color=TEXT_PRIMARY, bold=True, role="title")

# Left — Model bars
models = [
    ("Claude Opus 4.5", "~100K", "~180K", 0.5, 0.9),
    ("GPT-5.2", "~64K", "~200K", 0.32, 1.0),
    ("Gemini 3 Pro", "~500K", "~800K", 0.63, 1.0),
]
for i, (name, start, severe, s_pct, sev_pct) in enumerate(models):
    y = 240 + i * 200
    # Label
    add_text(s3, 96, y, 300, 40, name, font_size=20, color=TEXT_PRIMARY, bold=True, role="data")
    # Green bar (safe zone)
    bar_w = 800
    add_shape(s3, 96, y+50, int(bar_w * s_pct), 50, fill=STATE_CONVENTION)
    # Red bar (degradation zone)
    add_shape(s3, 96 + int(bar_w * s_pct), y+50, int(bar_w * (sev_pct - s_pct)), 50, fill=STATE_DECAY)
    # Labels
    add_text(s3, 96, y+110, 400, 30, f"\uc5f4\ud654\uc2dc\uc791 {start}", font_size=14, color=TEXT_SECONDARY)
    add_text(s3, 500, y+110, 400, 30, f"\uc2ec\uac01 {severe}", font_size=14, color=STATE_DECAY)

# Right — Evidence cards
evidence = [
    ("Attention Budget", "n\xb2 \uc5f0\uc0b0. \ud1a0\ud070\uc774 \ub298\uc218\ub85d \uac01 \ud1a0\ud070\uc5d0\n\ubc30\ubd84\ub418\ub294 \uc8fc\uc758\ub825\uc774 \ud76c\uc11d\ub41c\ub2e4.\n\ube44\uc6a9\uc774 \uc544\ub2c8\ub77c \ud488\uc9c8\uc758 \ubb38\uc81c.", STATE_OBSERVE),
    ("Lost-in-Middle", "U\uc790\ud615 \uc5b4\ud150\uc158 \uace1\uc120.\n\uc911\uac04 \uc815\ubcf4\uc758 recall\uc774 10-40% \ud558\ub77d.\n\uc2dc\uc791\uacfc \ub05d\ub9cc \uc798 \uae30\uc5b5\ud55c\ub2e4.", STATE_DECAY),
    ("\uc555\ucd95\uc758 \uad50\ud6c8", "\uac00\uc7a5 \ub9ce\uc774 \uc904\uc778 Opaque(99.3%)\uac00\n\uac00\uc7a5 \ub0ae\uc740 \ud488\uc9c8(3.35).\n\uc62c\ubc14\ub978 \uac83\uc744 \uc720\uc9c0\ud558\ub294 \uac8c \uc911\uc694.", STATE_APPROVAL),
]
for i, (title, body, color) in enumerate(evidence):
    add_card(s3, 1000, 210 + i*220, 824, 200, title, body, color)

# Bottom takeaway
add_text(s3, 96, 900, 1728, 100,
    "Window\uac00 200K\ub77c\ub3c4 200K\ub97c \ucc44\uc6b0\uba74 \uc548 \ub41c\ub2e4. \uc5f4\ud654 \uc2dc\uc791\uc810 \uc774\uc804\uc774 \uc2e4\uc9c8\uc801 \uc720\ud6a8 \ubc94\uc704\ub2e4.",
    font_size=20, color=TEXT_SECONDARY, alignment=PP_ALIGN.CENTER)


# ═══════════════════════════════════════════════════════════════
# SLIDE 4 — Insight: 반창고가 아니라 상처를 치료
# ═══════════════════════════════════════════════════════════════
s4 = prs.slides.add_slide(blank_layout)
add_bg(s4)
add_shape(s4, 0, 0, 1920, 1080, fill=BG_BASE)

add_text(s4, 96, 72, 1728, 90,
    "\ubc18\ucc3d\uace0\uac00 \uc544\ub2c8\ub77c \uc0c1\ucc98\ub97c \uce58\ub8cc\ud574\uc57c \ud55c\ub2e4",
    font_size=34, color=TEXT_PRIMARY, bold=True, role="title")

# Left loop — bad cycle
add_card(s4, 120, 250, 700, 500,
    "\ub098\uc05c \uad6c\uc870 + \ubb38\uc11c \ubcf4\uc0c1",
    "\ub098\uc05c \ud30c\uc77c \uad6c\uc870\n      \u2193\nClaude\uac00 \ubabb \ucc3e\uc74c\n      \u2193\nCLAUDE.md\uc5d0 \uc9c0\ub3c4 \ucd94\uac00\n      \u2193\n\ub9ac\ud329\ud1a0\ub9c1 \ud6c4 \uc9c0\ub3c4 \ubbf8\uac31\uc2e0\n      \u2193\nContext Poisoning\n      \u2193\n\ub2e4\uc2dc \ucc98\uc74c\uc73c\ub85c \u2026",
    STATE_DECAY)

# Center VS
add_text(s4, 860, 430, 200, 140, "VS", font_size=46, color=LINE_DEFAULT, bold=True, alignment=PP_ALIGN.CENTER, role="title")

# Right loop — good cycle
add_card(s4, 1100, 250, 700, 500,
    "\uc88b\uc740 \uad6c\uc870 + \ucd5c\uc18c \ubb38\uc11c",
    "\uc790\uba85\ud55c \ud30c\uc77c \uad6c\uc870\n      \u2193\nClaude\uac00 ls\ub85c \ubc1c\uacac\n      \u2193\nCLAUDE.md\uc5d0\ub294 \uc554\ubb35\uc9c0\ub9cc\n      \u2193\n\uad6c\uc870 = \ubb38\uc11c (\ub3d9\uae30\ud654 \ubd88\ud544\uc694)\n      \u2193\nPoisoning \uc5c6\uc74c\n      \u2193\n\uc548\uc815 \uc720\uc9c0",
    STATE_CONVENTION)

# Thesis
add_text(s4, 96, 820, 1728, 120,
    "\ud50c\ub7ec\uadf8\uc778\uc758 \uc5ed\ud560\uc740 CLAUDE.md\uc5d0 \ub354 \ub9ce\uc740 \uc815\ubcf4\ub97c \ucc44\uc6cc\ub123\ub294 \uac83\uc774 \uc544\ub2c8\ub77c,\n\ud504\ub85c\uc81d\ud2b8 \uc790\uccb4\ub97c Claude\uac00 \ud0d0\uc0c9\ud558\uae30 \uc26c\uc6b4 \uad6c\uc870\ub85c \ub9cc\ub4dc\ub294 \uac83\uc774\ub2e4.",
    font_size=22, color=STATE_OBSERVE, bold=True, alignment=PP_ALIGN.CENTER, role="title")


# ═══════════════════════════════════════════════════════════════
# SLIDE 5 — 넣으면 안 되는 것 vs 넣어야 하는 것 (암묵지)
# ═══════════════════════════════════════════════════════════════
s4b = prs.slides.add_slide(blank_layout)
add_bg(s4b)
add_shape(s4b, 0, 0, 1920, 1080, fill=BG_BASE)

add_text(s4b, 96, 72, 1728, 90,
    "CLAUDE.md \xb7 rules/ \uc5d0 \ubb58 \ub123\uace0 \ubb58 \ube7c\ub294\uac00",
    font_size=34, color=TEXT_PRIMARY, bold=True, role="title")

add_text(s4b, 96, 170, 1728, 50,
    "Static Context\ub294 \ub9e4 \uc138\uc158, \ub9e4 \ud134\ub9c8\ub2e4 \ub85c\ub4dc\ub41c\ub2e4. \ubd88\ud544\uc694\ud55c \uc815\ubcf4\ub294 \uc7a1\uc74c\uc774 \ub418\uace0, \ud544\uc218 \uc815\ubcf4\uc758 \ubd80\uc7ac\ub294 \uc2e4\uc218\uac00 \ub41c\ub2e4.",
    font_size=18, color=TEXT_SECONDARY, alignment=PP_ALIGN.CENTER)

# ── Left: 넣으면 안 되는 것 ──
add_shape(s4b, 96, 240, 860, 56, fill=RGBColor(0x2A, 0x15, 0x1E))
add_text(s4b, 120, 248, 812, 40,
    "\u274c  \ub123\uc73c\uba74 \uc548 \ub418\ub294 \uc815\ubcf4",
    font_size=22, color=STATE_DECAY, bold=True, role="title")

dont_items = [
    (
        "\uc544\ud0a4\ud14d\ucc98 \uc9c0\ub3c4",
        "ls src/ \ub85c \ubc1c\uacac \uac00\ub2a5",
        "\ud30c\uc77c\uc2dc\uc2a4\ud15c \uad6c\uc870\uc758 \ubcf5\uc81c\ubcf8.\n\ub9ac\ud329\ud1a0\ub9c1 \ud6c4 \ubbf8\uac31\uc2e0 \uc2dc Context Poisoning.",
    ),
    (
        "API \ubb38\uc11c / \ud0c0\uc785 \uc815\uc758",
        "\ucf54\ub4dc \uc790\uccb4\uac00 \uc815\uc758",
        "Read\ub85c \uc77d\uc73c\uba74 \uc815\ud655\ud55c \ucd5c\uc2e0 \uc815\ubcf4.\n\ubcf5\uc81c\ubcf8\uc740 \ub0a1\uc740 \uc815\ubcf4\ub85c \uc624\uc5fc.",
    ),
    (
        "\ud30c\uc77c\ubcc4 \uc124\uba85",
        "\ud30c\uc77c\uc774 \uc2a4\uc2a4\ub85c \uc124\uba85\ud574\uc57c \ud568",
        "\ud30c\uc77c\uba85\uc774 \uc790\uba85\ud558\uc9c0 \uc54a\uc73c\uba74 \uc124\uba85\uc774 \uc544\ub2c8\ub77c\n\uc774\ub984\uc744 \uace0\uccd0\uc57c \ud568. \ubb38\uc11c\ub294 \ubc18\ucc3d\uace0.",
    ),
    (
        "\uc77c\ubc18\uc801 \ubca0\uc2a4\ud2b8 \ud504\ub799\ud2f0\uc2a4",
        "Claude\uac00 \uc774\ubbf8 \uc54c\uace0 \uc788\uc74c",
        "\"\ucf54\ub4dc\ub97c \ub9ac\ub354\ube14\ud558\uac8c \uc4f0\uc138\uc694\" \uac19\uc740 \uc9c0\uc2dc\ub294\n\ubaa8\ub378 \ud559\uc2b5 \ub370\uc774\ud130\uc5d0 \uc774\ubbf8 \uc788\uc74c. \ud1a0\ud070 \ub0ad\ube44.",
    ),
]

for i, (title, reason, detail) in enumerate(dont_items):
    y = 310 + i * 175
    add_shape(s4b, 96, y, 860, 160, fill=BG_SURFACE1, border_color=STATE_DECAY)
    add_text(s4b, 120, y + 12, 812, 32, title, font_size=18, color=STATE_DECAY, bold=True, role="title")
    # reason badge
    add_shape(s4b, 120, y + 48, 300, 28, fill=RGBColor(0x2A, 0x15, 0x1E), border_color=STATE_DECAY)
    add_text(s4b, 128, y + 50, 284, 24, reason, font_size=12, color=STATE_DECAY, role="data")
    # detail
    add_text(s4b, 120, y + 84, 812, 64, detail, font_size=14, color=TEXT_SECONDARY)

# ── Right: 넣어야 하는 것 (암묵지) ──
add_shape(s4b, 964, 240, 860, 56, fill=RGBColor(0x0E, 0x25, 0x18))
add_text(s4b, 988, 248, 812, 40,
    "\u2705  \ub123\uc5b4\uc57c \ud558\ub294 \uc815\ubcf4 \u2014 \uc554\ubb35\uc9c0 (Implicit Knowledge)",
    font_size=22, color=STATE_CONVENTION, bold=True, role="title")

do_items = [
    (
        "\ucee8\ubca4\uc158",
        "\ucf54\ub4dc\uc5d0 \"\uc65c\"\ub294 \uc548 \uc801\ud600 \uc788\uc74c",
        "\"\uc5d0\ub7ec \ucc98\ub9ac\ub294 Result \ud0c0\uc785, try-catch \uc544\ub2d8\"\n10\uac1c \ud30c\uc77c\uc744 \uc77d\uc5b4\ub3c4 \"try-catch \uc548 \uc4f4\ub2e4\" \ub294 \ubcf4\uc774\uc9c0\ub9cc\n\"\uc4f0\uba74 \uc548 \ub41c\ub2e4\"\ub294 \ubcf4\uc774\uc9c0 \uc54a\ub294\ub2e4.",
    ),
    (
        "\uae08\uc9c0 \uaddc\uce59",
        "\ucf54\ub4dc\uc5d0 \"\uc5c6\ub294 \uac83\"\uc740 \ubc1c\uacac \ubd88\uac00",
        "\"any \ud0c0\uc785 \uae08\uc9c0\", \"ORM X \uc0ac\uc6a9 \uae08\uc9c0\"\n\uc5c6\ub294 \uac83\uc740 grep\uc73c\ub85c \ucc3e\uc744 \uc218 \uc5c6\ub2e4.\n\uba85\uc2dc\uc801 \uae08\uc9c0\ub9cc\uc774 \uc608\ubc29\ud560 \uc218 \uc788\ub2e4.",
    ),
    (
        "\ube44\uc790\uba85\ud55c \uc2e4\ud589 \ubc29\ubc95",
        "package.json\ub9cc\uc73c\ub85c \ucd94\uce21 \ubd88\uac00",
        "\"bun test --filter=unit\"\nClaude\uac00 bun test \ub85c \uc2e4\ud328 \u2192 \uc7ac\uc2dc\ub3c4 \u2192 \ub610 \uc2e4\ud328.\n\ud55c \uc904 \uba85\uc2dc\uba74 \ub9e4 \uc138\uc158 \uc808\uc57d.",
    ),
    (
        "\ube44\uc790\uba85\ud55c \uad00\uacc4",
        "import \uadf8\ub798\ud504\ub85c \uc548 \ubcf4\uc774\ub294 \uc758\uc874\uc131",
        "\"Service A\ub294 \ubc18\ub4dc\uc2dc Service B \ucd08\uae30\ud654 \ud6c4 \uc2dc\uc791\"\n\ub7f0\ud0c0\uc784 \uc21c\uc11c, \ud658\uacbd\ubcc0\uc218 \uc758\uc874 \ub4f1\n\ucf54\ub4dc \uad6c\uc870\uc5d0 \ub4dc\ub7ec\ub098\uc9c0 \uc54a\ub294 \uc81c\uc57d.",
    ),
]

for i, (title, reason, detail) in enumerate(do_items):
    y = 310 + i * 175
    add_shape(s4b, 964, y, 860, 160, fill=BG_SURFACE1, border_color=STATE_CONVENTION)
    add_text(s4b, 988, y + 12, 812, 32, title, font_size=18, color=STATE_CONVENTION, bold=True, role="title")
    # reason badge
    add_shape(s4b, 988, y + 48, 340, 28, fill=RGBColor(0x0E, 0x25, 0x18), border_color=STATE_CONVENTION)
    add_text(s4b, 996, y + 50, 324, 24, reason, font_size=12, color=STATE_CONVENTION, role="data")
    # detail
    add_text(s4b, 988, y + 84, 812, 64, detail, font_size=14, color=TEXT_SECONDARY)

# Bottom — 판별 기준
add_shape(s4b, 96, 1010, 1728, 50, fill=BG_SURFACE2)
add_text(s4b, 120, 1015, 1680, 40,
    "\ud310\ubcc4 \uae30\uc900: \ubc1c\uacac \ubd88\uac00\ub2a5 + \ub9e4 \uc138\uc158 \ud544\uc694 + \uc548\uc815\uc801 + \uace0\uc2e0\ud638 \u2014 4\uac00\uc9c0 \ubaa8\ub450 \ucda9\uc871\ud558\uba74 \uc554\ubb35\uc9c0, \ud558\ub098\ub77c\ub3c4 \ube60\uc9c0\uba74 \uc81c\uc678",
    font_size=16, color=STATE_OBSERVE, alignment=PP_ALIGN.CENTER, role="data")


# ═══════════════════════════════════════════════════════════════
# SLIDE 6 — Core Loop: Observe → Analyze → Act → Measure
# ═══════════════════════════════════════════════════════════════
s5 = prs.slides.add_slide(blank_layout)
add_bg(s5)
add_shape(s5, 0, 0, 1920, 1080, fill=BG_BASE)

add_text(s5, 96, 72, 1728, 90,
    "Auto Context \u2014 \uc790\ub3d9 \ucee8\ud14d\uc2a4\ud2b8 \ucd5c\uc801\ud654 \ub8e8\ud504",
    font_size=34, color=TEXT_PRIMARY, bold=True, role="title")

# Central cycle — 4 boxes in diamond
cycle_items = [
    (810, 200, "Observe", "\uc218\uc9d1", STATE_OBSERVE),    # top
    (1150, 420, "Analyze", "\ubd84\uc11d", STATE_CANDIDATE), # right
    (810, 640, "Act", "\uc801\uc6a9", STATE_CONVENTION),      # bottom
    (470, 420, "Measure", "\uce21\uc815", STATE_APPROVAL),    # left
]
for cx, cy, label, sub, color in cycle_items:
    add_shape(s5, cx, cy, 300, 140, fill=BG_SURFACE1, border_color=color, border_width=Pt(2))
    add_text(s5, cx+10, cy+15, 280, 60, label, font_size=28, color=color, bold=True, alignment=PP_ALIGN.CENTER, role="title")
    add_text(s5, cx+10, cy+80, 280, 40, sub, font_size=18, color=TEXT_SECONDARY, alignment=PP_ALIGN.CENTER)

# Arrows between (simple text arrows)
add_text(s5, 1050, 310, 100, 40, "\u2192", font_size=32, color=LINE_DEFAULT, alignment=PP_ALIGN.CENTER)
add_text(s5, 1050, 590, 100, 40, "\u2190", font_size=32, color=LINE_DEFAULT, alignment=PP_ALIGN.CENTER)
add_text(s5, 700, 260, 100, 40, "\u2193", font_size=32, color=LINE_DEFAULT, alignment=PP_ALIGN.CENTER)
add_text(s5, 700, 680, 100, 40, "\u2191", font_size=32, color=LINE_DEFAULT, alignment=PP_ALIGN.CENTER)

# Bottom 4 captions
captions = [
    (96, 840, "Observe", "Hooks\uac00 \ub9e4 \ub3c4\uad6c \uc0ac\uc6a9\ub9c8\ub2e4\nRAW \ub370\uc774\ud130\ub97c \uc218\uc9d1\n\u2192 SQLite\uc5d0 \uc800\uc7a5", STATE_OBSERVE),
    (546, 840, "Analyze", "Worker\uac00 polling\ud558\uba70\n\ud328\ud134 \ubd84\uc11d, \ub300\ud654 \uc694\uc57d\n\u2192 5\ucc28\uc6d0 \uc810\uc218 \uc0b0\ucd9c", STATE_CANDIDATE),
    (996, 840, "Act", "\uc554\ubb35\uc9c0 \u2192 rules/ \uc790\ub3d9 \uc0dd\uc131\n\uad6c\uc870 \ubb38\uc81c \u2192 offers/ \uc81c\uc548\nCLAUDE.md \ucd5c\uc18c \uac31\uc2e0", STATE_CONVENTION),
    (1446, 840, "Measure", "Navigability, Readability\nPredictability, Self-doc\nIsolation \u2014 5\ucc28\uc6d0 \ucd94\uc801", STATE_APPROVAL),
]
for x, y, title, body, color in captions:
    add_shape(s5, x, y, 420, 160, fill=BG_SURFACE1, border_color=color)
    add_text(s5, x+16, y+10, 388, 30, title, font_size=16, color=color, bold=True, role="title")
    add_text(s5, x+16, y+45, 388, 100, body, font_size=13, color=TEXT_SECONDARY)


# ═══════════════════════════════════════════════════════════════
# SLIDE 6 — Feedback Loop: 4가지 신호
# ═══════════════════════════════════════════════════════════════
s6 = prs.slides.add_slide(blank_layout)
add_bg(s6)
add_shape(s6, 0, 0, 1920, 1080, fill=BG_BASE)

add_text(s6, 96, 72, 1728, 90,
    "\uc138\uc158\uc5d0\uc11c \ubc30\uc6b4\ub2e4 \u2014 4\uac00\uc9c0 \uc2e0\ud638",
    font_size=34, color=TEXT_PRIMARY, bold=True, role="title")

signals = [
    (96, 220, "\uac80\uc0c9 \uc2e0\ud638 (Glob \xb7 Grep)", "\ubaa9\ud45c \ud30c\uc77c\uae4c\uc9c0 \ud3c9\uade0 \uac80\uc0c9 \ud69f\uc218 \u2192 Navigability\n\uc608: \"auth \uad00\ub828 \ud30c\uc77c\uc744 \ucc3e\ub294 \ub370 \ud3c9\uade0 4.2\ud68c \uc2dc\ub3c4\"", STATE_OBSERVE),
    (984, 220, "\uc77d\uae30 \uc2e0\ud638 (Read)", "Read \uc904 \uc218 vs Edit \uc904 \uc218 \ube44\uc728 \u2192 Readability\n\uc608: \"utils.ts 245\uc904 Read, 10\uc904 Edit \u2192 \uc2e0\ud638 4%\"", STATE_OBSERVE),
    (96, 430, "\uc218\uc815 \uc2e0\ud638 (Edit)", "\ubc18\ubcf5\ub418\ub294 \ubcc0\ud658 \ud328\ud134 \u2192 \ubbf8\uba85\ubb38\ud654\ub41c \ucee8\ubca4\uc158\n\uc608: \"try-catch \u2192 Result \ubcc0\ud658 5\uc138\uc158 \uc5f0\uc18d\"", STATE_CONVENTION),
    (984, 430, "\uc2e4\ud589 \uc2e0\ud638 (Bash)", "\uc2e4\ud328 \ud6c4 \uc131\uacf5\ud55c \uba85\ub839\uc5b4 \ud328\ud134 \u2192 \ube44\uc790\uba85\ud55c \uc2e4\ud589 \ubc29\ubc95\n\uc608: \"bun test \uc2e4\ud328 \u2192 bun test --filter \uc131\uacf5\"", STATE_APPROVAL),
]
for x, y, title, body, color in signals:
    add_card(s6, x, y, 840, 170, title, body, color)

# Summary bar
add_shape(s6, 96, 650, 1728, 180, fill=BG_SURFACE2, border_color=STATE_CANDIDATE)
add_text(s6, 140, 670, 1640, 140,
    "Hook\uc740 \ubc14\ubcf4 \uc218\uc9d1\uae30 \u2014 \ubc1b\uc544\uc11c \ub358\uc9c0\uae30\ub9cc.\n\ubaa8\ub4e0 \ubd84\uc11d\uc740 Background Worker\uc5d0\uc11c.\nMain Session\uc5d0 \ub808\uc774\ud134\uc2dc 0.",
    font_size=20, color=TEXT_SECONDARY, alignment=PP_ALIGN.CENTER)


# ═══════════════════════════════════════════════════════════════
# SLIDE 7 — Safety: 자동 vs 승인
# ═══════════════════════════════════════════════════════════════
s7 = prs.slides.add_slide(blank_layout)
add_bg(s7)
add_shape(s7, 0, 0, 1920, 1080, fill=BG_BASE)

add_text(s7, 96, 72, 1728, 90,
    "\uc790\ub3d9\ud654\uc758 \uacbd\uacc4 \u2014 \ubb34\uc5c7\uc774 \uc790\ub3d9\uc774\uace0 \ubb34\uc5c7\uc774 \uc2b9\uc778\uc778\uac00",
    font_size=34, color=TEXT_PRIMARY, bold=True, role="title")

add_text(s7, 96, 190, 1728, 80,
    "Worker\ub294 \ud504\ub85c\uc81d\ud2b8\ub97c \uad00\ucc30\ud558\uace0 \ud559\uc2b5\ud558\uc9c0\ub9cc, \uad6c\uc870 \ubcc0\uacbd\uc740 \uc0ac\uc6a9\uc790 \uc2b9\uc778 \uc5c6\uc774 \ud558\uc9c0 \uc54a\ub294\ub2e4.",
    font_size=22, color=TEXT_SECONDARY, alignment=PP_ALIGN.CENTER)

# Safety table
safety_data = [
    ["\uc0b0\ucd9c\ubb3c", "\uc790\ub3d9", "\uc2b9\uc778", "\uc124\uba85"],
    [".claude/rules/ \uc0dd\uc131", "\u2713", "", "\ubc18\ubcf5 \ud328\ud134\uc5d0\uc11c \uc554\ubb35\uc9c0 \ucd94\ucd9c. \ucf54\ub4dc \ubcc0\uacbd \uc544\ub2d8."],
    ["CLAUDE.md \uac31\uc2e0", "\u2713", "", "\ube44\uc790\uba85\ud55c \uc2e4\ud589 \ubc29\ubc95 \ucd94\uac00. \ucf54\ub4dc \ubcc0\uacbd \uc544\ub2d8."],
    ["offers/ \uad6c\uc870 \uc81c\uc548", "\u2713", "", "Worker\uac00 \uc81c\uc548 \ud30c\uc77c \uc0dd\uc131. \uc544\uc9c1 \uc544\ubb34\uac83\ub3c4 \uc548 \ubc14\ub00c."],
    ["\uad6c\uc870 \ub9ac\ud329\ud1a0\ub9c1 \uc2e4\ud589", "", "\u2713", "/cac-apply\ub85c \uc0ac\uc6a9\uc790\uac00 \uc120\ud0dd\ud55c offer\ub9cc \uc2e4\ud589."],
]
add_table_slide(s7, 180, 300, 1560, safety_data, [0.22, 0.1, 0.1, 0.58])

add_text(s7, 96, 860, 1728, 100,
    "\ucf54\ub4dc\uac00 \ubc14\ub00c\ub294 \uac74 \uc624\uc9c1 /cac-apply\ub97c \uc2e4\ud589\ud560 \ub54c\ubfd0. \uadf8 \uc678\uc5d0\ub294 \uad00\ucc30\ud558\uace0 \uc81c\uc548\ud560 \ubfd0\uc774\ub2e4.",
    font_size=20, color=STATE_APPROVAL, bold=True, alignment=PP_ALIGN.CENTER, role="title")


# ═══════════════════════════════════════════════════════════════
# SLIDE 8 — Two-Speed Architecture
# ═══════════════════════════════════════════════════════════════
s8 = prs.slides.add_slide(blank_layout)
add_bg(s8)
add_shape(s8, 0, 0, 1920, 1080, fill=BG_BASE)

add_text(s8, 96, 72, 1728, 90,
    "Two-Speed \uc544\ud0a4\ud14d\ucc98 \u2014 \ube60\ub978 \uc218\uc9d1, \ub290\ub9b0 \ubd84\uc11d",
    font_size=34, color=TEXT_PRIMARY, bold=True, role="title")

# Left — Fast Lane
add_shape(s8, 96, 220, 820, 620, fill=BG_SURFACE1, border_color=STATE_OBSERVE)
add_text(s8, 120, 235, 780, 45, "Fast Lane \u2014 \uc218\uc9d1", font_size=24, color=STATE_OBSERVE, bold=True, role="title")
add_text(s8, 120, 290, 780, 530,
    "PostToolUse \uc774\ubca4\ud2b8 \ubc1c\uc0dd\n      \u2193\nHook (\uc178 \uc2a4\ud06c\ub9bd\ud2b8)\n      \ubc14\ubcf4. pipe\ub9cc.\n      \u2193\ncollector.mjs\n      JSON \ud30c\uc2f1\n      parameterized INSERT\n      \u2193\nSQLite raw_events\n      processed = 0\n\n\uc0dd\uba85\uc8fc\uae30: fire-and-forget\n\ub808\uc774\ud134\uc2dc: ms \ub2e8\uc704\nMain Session \uc601\ud5a5: \uc5c6\uc74c",
    font_size=15, color=TEXT_SECONDARY, role="data")

# Mid — SQLite bus
add_shape(s8, 940, 260, 40, 540, fill=LINE_DEFAULT)
add_text(s8, 920, 490, 80, 40, "SQLite", font_size=14, color=TEXT_PRIMARY, bold=True, alignment=PP_ALIGN.CENTER, role="data")

# Right — Slow Lane
add_shape(s8, 1004, 220, 820, 620, fill=BG_SURFACE1, border_color=STATE_CANDIDATE)
add_text(s8, 1028, 235, 780, 45, "Slow Lane \u2014 \ubd84\uc11d", font_size=24, color=STATE_CANDIDATE, bold=True, role="title")
add_text(s8, 1028, 290, 780, 530,
    "Worker (polling loop)\n      \u2193\nprocessed=0 \uc870\ud68c\n      \u2193\nTool Events \ubd84\uc11d\n      \u2192 5\ucc28\uc6d0 \uc810\uc218 \uc0b0\ucd9c\n      \u2193\n\ub300\ud654 \uc694\uc57d (Stop RAW)\n      \u2192 sessions \ud14c\uc774\ube14\n      \u2193\n\ub204\uc801 \ubd84\uc11d (N\uc138\uc158)\n      \u2192 rules/ \uc0dd\uc131\n      \u2192 offers/ \uc0dd\uc131\n\n\uc0dd\uba85\uc8fc\uae30: \uc0c1\uc8fc \ud504\ub85c\uc138\uc2a4\nMain Session\uacfc: \uc644\uc804 \ub3c5\ub9bd",
    font_size=15, color=TEXT_SECONDARY, role="data")

# Bottom strip
add_shape(s8, 96, 880, 1728, 100, fill=BG_SURFACE2, border_color=LINE_DEFAULT)
add_text(s8, 120, 895, 1680, 70,
    "Glob \u2192 Read \u2192 Edit \u2192 Read \u2192 Bash \u2192 Read \u2192 Edit \u2192 Stop\n\uc804\ubd80 RAW\ub85c SQLite\uc5d0. \uac00\uacf5 \uc5c6\uc74c. \ud310\ub2e8 \uc5c6\uc74c. Worker\uac00 \uc54c\uc544\uc11c.",
    font_size=16, color=TEXT_SECONDARY, alignment=PP_ALIGN.CENTER, role="data")


# ═══════════════════════════════════════════════════════════════
# SLIDE 9 — Tech Stack
# ═══════════════════════════════════════════════════════════════
s9 = prs.slides.add_slide(blank_layout)
add_bg(s9)
add_shape(s9, 0, 0, 1920, 1080, fill=BG_BASE)

add_text(s9, 96, 72, 1728, 90,
    "\uae30\uc220 \uc2a4\ud0dd",
    font_size=34, color=TEXT_PRIMARY, bold=True, role="title")

add_text(s9, 96, 180, 1728, 70,
    "\ucd5c\uc18c \uc758\uc874\uc131. Claude Code\uc758 Node.js \ud658\uacbd\ub9cc\uc73c\ub85c \ub3d9\uc791\ud55c\ub2e4.",
    font_size=22, color=TEXT_SECONDARY, alignment=PP_ALIGN.CENTER)

stack_data = [
    ["\ucef4\ud3ec\ub10c\ud2b8", "\uae30\uc220", "\uc124\uba85"],
    ["Hook", "Bash (\uc178 \uc2a4\ud06c\ub9bd\ud2b8)", "PostToolUse, Stop \uc774\ubca4\ud2b8\ub97c \ubc1b\uc544 collector\uc5d0 pipe"],
    ["Collector", "Node.js (collector.mjs)", "JSON \ud30c\uc2f1 + SQLite parameterized INSERT. \uc774\ubca4\ud2b8\ub2f9 1\ud68c \uc2e4\ud589 \ud6c4 \uc885\ub8cc"],
    ["Storage", "SQLite (better-sqlite3)", "raw_events \xb7 sessions \xb7 insights 3\uac1c \ud14c\uc774\ube14"],
    ["Worker", "Claude Code (\ubcc4\ub3c4 \ud504\ub85c\uc138\uc2a4)", "SQLite polling \u2192 RAW \ubd84\uc11d \u2192 \uc694\uc57d \xb7 \ud328\ud134 \ucd94\ucd9c \xb7 5\ucc28\uc6d0 \uc810\uc218"],
    ["Input", ".claude/rules/ (Markdown)", "glob \uc2a4\ucf54\ud551. \ud30c\uc77c \uc218\uc815 \uc2dc \uc790\ub3d9 \ub85c\ub4dc. Worker\uac00 \uc0dd\uc131/\uac31\uc2e0"],
    ["Output", ".claude-auto-context/offers/", "\uad6c\uc870 \ubcc0\uacbd \uc81c\uc548. /cac-apply \uc2a4\ud0ac\ub85c \uc2b9\uc778 \ud6c4 \uc2e4\ud589"],
    ["Notification", "UserPromptSubmit Hook", "pending offer \ud655\uc778 \u2192 \uc54c\ub9bc \uc8fc\uc785. offer \uc5c6\uc73c\uba74 \ud328\uc2a4"],
]
add_table_slide(s9, 96, 280, 1728, stack_data, [0.15, 0.25, 0.60])


# ═══════════════════════════════════════════════════════════════
# SLIDE 10 — UX Demo: 5단계 사용자 경험
# ═══════════════════════════════════════════════════════════════
s10 = prs.slides.add_slide(blank_layout)
add_bg(s10)
add_shape(s10, 0, 0, 1920, 1080, fill=BG_BASE)

add_text(s10, 96, 72, 1728, 90,
    "\uc0ac\uc6a9\uc790 \uacbd\ud5d8 \u2014 5\ub2e8\uacc4",
    font_size=34, color=TEXT_PRIMARY, bold=True, role="title")

steps = [
    ("\ud3c9\uc18c\ucc98\ub7fc \ucf54\ub529", "Claude Code\ub85c \uc77c\uc0c1\uc801 \uc791\uc5c5.\n\ud50c\ub7ec\uadf8\uc778\uc740 \ubcf4\uc774\uc9c0\n\uc54a\uac8c \ub3d9\uc791.", STATE_OBSERVE),
    ("\uc790\ub3d9 \uc218\uc9d1", "\ub9e4 \ub3c4\uad6c \uc0ac\uc6a9\ub9c8\ub2e4\nHook \u2192 collector\n\u2192 SQLite", STATE_OBSERVE),
    ("\ubc31\uadf8\ub77c\uc6b4\ub4dc \ubd84\uc11d", "Worker\uac00 \ud328\ud134 \uac10\uc9c0.\n5\ucc28\uc6d0 \uc810\uc218 \uc0b0\ucd9c.", STATE_CANDIDATE),
    ("\uc54c\ub9bc \uc218\uc2e0", "\ub2e4\uc74c \uc138\uc158\uc5d0\uc11c\n\"N\uac74\uc758 Offer \ub300\uae30 \uc911\"\n\uc54c\ub9bc", STATE_APPROVAL),
    ("/cac-apply", "offer \uc120\ud0dd \u2192 Claude\uac00\n\uc790\ub3d9 \ub9ac\ud329\ud1a0\ub9c1\n+ \ud14c\uc2a4\ud2b8 \ud655\uc778", STATE_CONVENTION),
]

step_w = 320
gap = (1728 - step_w * 5) // 4
for i, (title, body, color) in enumerate(steps):
    x = 96 + i * (step_w + gap)
    # Step number circle
    add_shape(s10, x + step_w//2 - 25, 210, 50, 50, fill=color)
    add_text(s10, x + step_w//2 - 25, 215, 50, 45, str(i+1), font_size=22, color=BG_BASE, bold=True, alignment=PP_ALIGN.CENTER, role="data")
    # Arrow (except last)
    if i < 4:
        arrow_x = x + step_w + gap//2 - 15
        add_text(s10, arrow_x, 215, 30, 45, "\u2192", font_size=24, color=LINE_DEFAULT, alignment=PP_ALIGN.CENTER)
    # Card
    add_shape(s10, x, 280, step_w, 340, fill=BG_SURFACE1, border_color=color)
    add_text(s10, x+16, 295, step_w-32, 40, title, font_size=18, color=color, bold=True, alignment=PP_ALIGN.CENTER, role="title")
    add_text(s10, x+16, 350, step_w-32, 250, body, font_size=14, color=TEXT_SECONDARY, alignment=PP_ALIGN.CENTER)

# Result panel
add_shape(s10, 96, 700, 1728, 250, fill=BG_SURFACE2, border_color=STATE_APPROVAL)
add_text(s10, 140, 720, 1640, 210,
    "\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\n\U0001f514 Auto Context \u2014 2\uac74\uc758 Offer \ub300\uae30 \uc911\n\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\n1. src/utils.ts \ubd84\ud560 (Readability +0.5)\n2. routes/ \ud328\ud134 \ud1b5\uc77c (Predictability +0.8)\n\U0001f4a1 /cac-apply \ub85c \uc801\uc6a9 \xb7 /cac-status \ub85c \uc0c1\uc138 \ud655\uc778",
    font_size=20, color=STATE_APPROVAL, role="data")


# ═══════════════════════════════════════════════════════════════
# SLIDE 11 — Summary: Before / After
# ═══════════════════════════════════════════════════════════════
s11 = prs.slides.add_slide(blank_layout)
add_bg(s11)
add_shape(s11, 0, 0, 1920, 1080, fill=BG_BASE)

add_text(s11, 96, 72, 1728, 90,
    "Auto Context\uac00 \ubc14\uafb8\ub294 \uac83",
    font_size=34, color=TEXT_PRIMARY, bold=True, role="title")

summary_data = [
    ["\ucc28\uc6d0", "Before", "After"],
    ["CLAUDE.md", "\uc544\ud0a4\ud14d\ucc98 \uc9c0\ub3c4, API \ubb38\uc11c, \ud30c\uc77c \uc124\uba85\u2026 \uc218\ubc31 \uc904", "\uc554\ubb35\uc9c0\ub9cc. \ucee8\ubca4\uc158, \uae08\uc9c0 \uaddc\uce59, \ube44\uc790\uba85\ud55c \uad00\uacc4. \ucd5c\uc18c\ud55c."],
    ["\ucee8\ubca4\uc158", "\uc0ac\ub78c\uc774 \uae30\uc5b5\uc5d0 \uc758\uc874. \uc138\uc158\ub9c8\ub2e4 \ubd88\uc77c\uce58.", "Worker\uac00 \uc790\ub3d9 \ucd94\ucd9c \u2192 rules/ \uba85\ubb38\ud654."],
    ["\ud30c\uc77c \uad6c\uc870", "utils.ts 800\uc904. Claude\uac00 \ub9e4\ubc88 \uc804\ubd80 \uc77d\uc74c.", "offer\ub85c \ubd84\ud560 \uc81c\uc548 \u2192 /cac-apply\ub85c \ub9ac\ud329\ud1a0\ub9c1."],
    ["\ud0d0\uc0c9 \ud6a8\uc728", "\ubaa9\ud45c \ud30c\uc77c\uae4c\uc9c0 \ud3c9\uade0 4+ \uac80\uc0c9.", "\uad6c\uc870 \uac1c\uc120 \ud6c4 1-2\ud68c\ub85c \uc218\ub834."],
    ["\ud559\uc2b5", "\uc5c6\uc74c. \ub9e4 \uc138\uc158\uc774 \uccab \ub0a0.", "SQLite\uc5d0 \uc138\uc158 \ucd95\uc801. Worker\uac00 \ub204\uc801 \ud559\uc2b5."],
    ["\uc548\uc804", "\uc790\ub3d9\ud654 \uc5c6\uc74c. \uc804\ubd80 \uc218\ub3d9.", "\uad00\ucc30\uc740 \uc790\ub3d9, \uad6c\uc870 \ubcc0\uacbd\uc740 \uc2b9\uc778 \ud6c4\ub9cc."],
]
tbl = add_table_slide(s11, 96, 220, 1728, summary_data, [0.14, 0.43, 0.43])

# Color Before column red, After column green
table = tbl.table
for ri in range(1, len(summary_data)):
    # Before cell text color
    for p in table.cell(ri, 1).text_frame.paragraphs:
        p.font.color.rgb = STATE_DECAY
    # After cell text color
    for p in table.cell(ri, 2).text_frame.paragraphs:
        p.font.color.rgb = STATE_CONVENTION

# Closing
add_text(s11, 96, 880, 1728, 120,
    "\ud504\ub85c\uc81d\ud2b8\uc5d0 \ub354 \ub9ce\uc740 \ubb38\uc11c\ub97c \ub123\ub294 \uac83\uc774 \uc544\ub2c8\ub77c,\n\ud504\ub85c\uc81d\ud2b8 \uc790\uccb4\uac00 \ubb38\uc11c\uac00 \ub418\uac8c \ub9cc\ub4dc\ub294 \uac83. Auto Context.",
    font_size=24, color=TEXT_PRIMARY, bold=True, alignment=PP_ALIGN.CENTER, role="title")


# ═══════════════════════════════════════════════════════════════
# SLIDE 13 — Evolution: auto-context → claude-auto-context
# ═══════════════════════════════════════════════════════════════
s13 = prs.slides.add_slide(blank_layout)
add_bg(s13)
add_shape(s13, 0, 0, 1920, 1080, fill=BG_BASE)

add_text(s13, 96, 72, 1728, 90,
    "auto-context \u2192 claude-auto-context",
    font_size=46, color=TEXT_PRIMARY, bold=True, role="title")

add_text(s13, 96, 170, 1728, 50,
    "\uac19\uc740 \ubb38\uc81c\ub97c \ud480\uc9c0\ub9cc, \uc124\uacc4 \ucca0\ud559\uc774 \ub2e4\ub974\ub2e4.",
    font_size=24, color=TEXT_SECONDARY, alignment=PP_ALIGN.CENTER)

# ── Left column header: auto-context (어제) ──
add_shape(s13, 96, 240, 860, 52, fill=RGBColor(0x2A, 0x15, 0x1E))
add_text(s13, 120, 248, 400, 36,
    "auto-context", font_size=22, color=STATE_DECAY, bold=True, role="data")
add_text(s13, 540, 250, 380, 36,
    "\ub2e8\uc77c \ud30c\uc774\ud504\ub77c\uc778 \xb7 zero-dep",
    font_size=14, color=STATE_DECAY, role="data")

# ── Right column header: claude-auto-context (오늘) ──
add_shape(s13, 964, 240, 860, 52, fill=RGBColor(0x0E, 0x25, 0x18))
add_text(s13, 988, 248, 440, 36,
    "claude-auto-context", font_size=22, color=STATE_CONVENTION, bold=True, role="data")
add_text(s13, 1440, 250, 360, 36,
    "Two-Speed \xb7 \uad6c\uc870 \uce58\ub8cc",
    font_size=14, color=STATE_CONVENTION, role="data")

# ── Comparison rows ──
rows = [
    (
        "\uc800\uc7a5\uc18c",
        "*.json / .jsonl \ud30c\uc77c",
        "SQLite (raw_events \xb7 sessions \xb7 insights)",
    ),
    (
        "\uc544\ud0a4\ud14d\ucc98",
        "Hook + Stop agent + SessionStart\n\uc2a4\ud06c\ub9bd\ud2b8\uac00 \uc9c1\uc811 \ucc98\ub9ac\n(\ub2e8\uc77c \ud30c\uc774\ud504\ub77c\uc778)",
        "Hook\uc740 \uc218\uc9d1\ub9cc, Worker\uac00 polling \ubd84\uc11d\n(Two-Speed \u2014 \ube60\ub978 \uc218\uc9d1 / \ub290\ub9b0 \ubd84\uc11d)",
    ),
    (
        "\uc218\uc9d1 \uc2e0\ud638",
        "Write / Edit / Bash \uc2e4\ud328\n+ \uba85\uc2dc \ud53c\ub4dc\ubc31",
        "Glob \xb7 Grep \xb7 Read \xb7 Edit \xb7 Bash\n5\ucc28\uc6d0 \uc815\ub7c9\ud654",
    ),
    (
        "\ud575\uc2ec \ubaa9\ud45c",
        "\ucee8\ubca4\uc158 / \uc548\ud2f0\ud328\ud134 \ud559\uc2b5 \ud6c4\nCLAUDE.md + rules \uc7ac\uc0dd\uc131",
        "\ud504\ub85c\uc81d\ud2b8 \uad6c\uc870 \uc790\uccb4 \uac1c\uc120\n(offers \u2192 /cac-apply)",
    ),
    (
        "\uc2b9\uc778 \uacbd\uacc4",
        "/ac-review \uc2b9\uc778 \uac8c\uc774\ud2b8\n\uba85\uc2dc \ud53c\ub4dc\ubc31\uc740 \uc989\uc2dc active",
        "\uad6c\uc870 \ubcc0\uacbd\ub9cc /cac-apply \uc2b9\uc778\nrules \xb7 CLAUDE.md\ub294 \uc790\ub3d9",
    ),
    (
        "\ucd9c\ub825\ubb3c",
        "CLAUDE.md \ub9c8\ucee4 \uad6c\uac04\n+ auto-context-*.md",
        ".claude/rules/ \uc790\ub3d9 \uc0dd\uc131\n+ offers/ \uad6c\uc870 \uc81c\uc548",
    ),
    (
        "\ub7f0\ud0c0\uc784",
        "Bash + jq + JSON\nzero-dependency",
        "Node.js collector + SQLite\n+ background Claude worker",
    ),
]

row_h = 96
start_y = 308
for i, (label, left_text, right_text) in enumerate(rows):
    y = start_y + i * (row_h + 8)
    # Label
    add_shape(s13, 96, y, 860, row_h, fill=BG_SURFACE1, border_color=STATE_DECAY)
    add_shape(s13, 964, y, 860, row_h, fill=BG_SURFACE1, border_color=STATE_CONVENTION)

    # Dimension label (left-aligned inside left card)
    add_text(s13, 112, y + 6, 180, 24, label,
        font_size=13, color=STATE_OBSERVE, bold=True, role="data")
    # Left content
    add_text(s13, 112, y + 32, 820, row_h - 38, left_text,
        font_size=13, color=TEXT_SECONDARY)
    # Right content
    add_text(s13, 980, y + 6, 180, 24, label,
        font_size=13, color=STATE_OBSERVE, bold=True, role="data")
    add_text(s13, 980, y + 32, 820, row_h - 38, right_text,
        font_size=13, color=TEXT_SECONDARY)

# Bottom — key difference
add_shape(s13, 96, 1020, 1728, 44, fill=BG_SURFACE2)
add_text(s13, 120, 1026, 1680, 32,
    "auto-context: \ubb38\uc11c\ub97c \ub354 \uc798 \uc4f0\uc790   \u2192   claude-auto-context: \ubb38\uc11c\uac00 \ud544\uc694 \uc5c6\uac8c \uad6c\uc870\ub97c \uace0\uce58\uc790",
    font_size=18, color=STATE_CANDIDATE, bold=True, alignment=PP_ALIGN.CENTER, role="title")


# ─── Save ─────────────────────────────────────────────────────
output_path = os.path.join(os.path.dirname(os.path.dirname(__file__)), "Auto-Context.pptx")
prs.save(output_path)
print(f"PPT saved to: {output_path}")
